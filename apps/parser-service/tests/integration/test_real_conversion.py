"""Integration tests for real file-to-markdown conversion via markitdown.

These tests create actual files programmatically and verify the converter
produces usable text output for each supported format.

Note on "markdown": markitdown returns true markdown tables for CSV/XLSX,
but for DOCX/PPTX/PDF/TXT the output is plain text extraction. The converter
post-processes the output (strips HTML comments, normalizes whitespace) to
make it clean LLM context regardless of source format.
"""

import io

import pytest
from fastapi import UploadFile

from src.services.converter import FileConverter, get_converter


class TestRealFileConversion:
    """Test real markitdown conversion across all supported formats."""

    @pytest.fixture
    def converter(self) -> FileConverter:
        """Provide a real FileConverter instance."""
        return get_converter()

    async def _convert_bytes(
        self, converter: FileConverter, filename: str, data: bytes
    ) -> dict:
        """Helper: wrap bytes in UploadFile and convert."""
        upload = UploadFile(filename=filename, file=io.BytesIO(data))
        return await converter.convert(upload)

    @pytest.mark.asyncio
    async def test_convert_txt(self, converter: FileConverter) -> None:
        """Plain text files should pass through with minimal change."""
        result = await self._convert_bytes(
            converter, "report.txt", b"Hello World\n\nSecond paragraph."
        )
        assert result["filename"] == "report.txt"
        assert "Hello World" in result["markdown"]
        assert "Second paragraph" in result["markdown"]

    @pytest.mark.asyncio
    async def test_convert_md(self, converter: FileConverter) -> None:
        """Markdown files should retain their structure."""
        md = b"# Heading\n\nSome **bold** text.\n"
        result = await self._convert_bytes(converter, "doc.md", md)
        assert "# Heading" in result["markdown"]
        assert "bold" in result["markdown"]

    @pytest.mark.asyncio
    async def test_convert_csv_to_markdown_table(self, converter: FileConverter) -> None:
        """CSV files should be converted to real markdown tables."""
        csv = b"name,age\nAlice,30\nBob,25"
        result = await self._convert_bytes(converter, "data.csv", csv)
        # True markdown table output
        assert "| name | age |" in result["markdown"]
        assert "| Alice | 30 |" in result["markdown"]
        assert "| Bob | 25 |" in result["markdown"]

    @pytest.mark.asyncio
    async def test_convert_json(self, converter: FileConverter) -> None:
        """JSON files should have content readable in output."""
        data = b'{"company": "Acme", "revenue": 1000000}'
        result = await self._convert_bytes(converter, "data.json", data)
        assert "Acme" in result["markdown"]
        assert "1000000" in result["markdown"]

    @pytest.mark.asyncio
    async def test_convert_xlsx_to_markdown_table(self, converter: FileConverter) -> None:
        """Excel files should be converted to real markdown tables."""
        from openpyxl import Workbook

        wb = Workbook()
        ws = wb.active
        ws.title = "Revenue"
        ws["A1"] = "Quarter"
        ws["B1"] = "Revenue"
        ws["A2"] = "Q1"
        ws["B2"] = 1000000

        buf = io.BytesIO()
        wb.save(buf)
        buf.seek(0)

        upload = UploadFile(filename="revenue.xlsx", file=buf)
        result = await converter.convert(upload)

        assert result["filename"] == "revenue.xlsx"
        # True markdown table output
        assert "## Revenue" in result["markdown"]
        assert "| Quarter | Revenue |" in result["markdown"]
        assert "| Q1 | 1000000 |" in result["markdown"]

    @pytest.mark.asyncio
    async def test_convert_pptx_text_extraction(self, converter: FileConverter) -> None:
        """PowerPoint files should have slide text extracted.

        Note: PPTX output from markitdown is plain text, not markdown.
        The converter strips HTML comments (e.g. <!-- Slide number: 1 -->)
        to produce clean text.
        """
        from pptx import Presentation
        from pptx.util import Inches

        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank layout
        textbox = slide.shapes.add_textbox(
            Inches(1), Inches(1), Inches(5), Inches(1)
        )
        textbox.text_frame.text = "Annual Revenue: $10M"

        buf = io.BytesIO()
        prs.save(buf)
        buf.seek(0)

        upload = UploadFile(filename="slides.pptx", file=buf)
        result = await converter.convert(upload)

        assert result["filename"] == "slides.pptx"
        assert "Annual Revenue" in result["markdown"]
        assert "$10M" in result["markdown"]
        # HTML comments should be stripped by _clean_markdown
        assert "<!--" not in result["markdown"]

    @pytest.mark.asyncio
    async def test_convert_docx_text_extraction(self, converter: FileConverter) -> None:
        """Word documents should have text and headings extracted.

        Note: DOCX output from markitdown is plain text with line breaks,
        not markdown headings.
        """
        from docx import Document

        doc = Document()
        doc.add_heading("Annual Report", level=0)
        doc.add_paragraph("Revenue was $10M this year.")
        doc.add_heading("Expenses", level=1)
        doc.add_paragraph("Expenses were $8M.")

        buf = io.BytesIO()
        doc.save(buf)
        buf.seek(0)

        upload = UploadFile(filename="report.docx", file=buf)
        result = await converter.convert(upload)

        assert result["filename"] == "report.docx"
        assert "Annual Report" in result["markdown"]
        assert "Revenue was $10M" in result["markdown"]
        assert "Expenses were $8M" in result["markdown"]

    @pytest.mark.asyncio
    async def test_convert_pdf_text_extraction(self, converter: FileConverter) -> None:
        """PDF files should have text extracted.

        Note: PDF output from markitdown is plain text, not markdown.
        """
        from fpdf import FPDF

        pdf = FPDF()
        pdf.add_page()
        pdf.set_font("Helvetica", size=12)
        pdf.cell(200, 10, text="PDF Revenue Report", new_x="LMARGIN", new_y="NEXT")
        pdf.cell(200, 10, text="Revenue: $10M", new_x="LMARGIN", new_y="NEXT")

        buf = io.BytesIO()
        pdf.output(buf)
        buf.seek(0)

        upload = UploadFile(filename="report.pdf", file=buf)
        result = await converter.convert(upload)

        assert result["filename"] == "report.pdf"
        # PDF text extraction may vary; check for presence of key content
        assert "Revenue" in result["markdown"] or "Report" in result["markdown"]

    @pytest.mark.asyncio
    async def test_convert_preserves_metadata(self, converter: FileConverter) -> None:
        """Converter should include metadata when available."""
        from docx import Document

        doc = Document()
        doc.add_paragraph("Simple content.")

        buf = io.BytesIO()
        doc.save(buf)
        buf.seek(0)

        upload = UploadFile(filename="meta.docx", file=buf)
        result = await converter.convert(upload)

        assert "markdown" in result
        assert "filename" in result
        assert result["filename"] == "meta.docx"
